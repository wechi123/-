# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

P=Presentation();P.slide_width=Inches(13.333);P.slide_height=Inches(7.5)
D=RGBColor(0x1A,0x3C,0x6E);A=RGBColor(0x2B,0x57,0x9A);L=RGBColor(0x3A,0x7C,0xD5)
W=RGBColor(0xFF,0xFF,0xFF);G=RGBColor(0xF0,0xF2,0xF5);K=RGBColor(0x33,0x33,0x33)
O=RGBColor(0xD4,0xA0,0x1F);R=RGBColor(0x2E,0x7D,0x32);Y=RGBColor(0xAA,0xBB,0xCC)

def Bx(s,l,t,w,h,c,st=MSO_SHAPE.RECTANGLE):
 x=s.shapes.add_shape(st,l,t,w,h);x.fill.solid();x.fill.fore_color.rgb=c;x.line.fill.background()
def Tx(s,l,t,w,h,tx,fs=18,c=K,b=False,a=PP_ALIGN.LEFT):
 x=s.shapes.add_textbox(l,t,w,h);f=x.text_frame;f.word_wrap=True
 p=f.paragraphs[0];p.text=tx;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.bold=b
 p.font.name='Microsoft YaHei';p.alignment=a
def Li(s,l,t,w,h,items,fs=16,c=K,sp=Pt(8)):
 x=s.shapes.add_textbox(l,t,w,h);f=x.text_frame;f.word_wrap=True
 for i,item in enumerate(items):
  p=f.paragraphs[0] if i==0 else f.add_paragraph()
  p.text=item;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.name='Microsoft YaHei';p.space_after=sp
def TB(s,tx):Bx(s,Inches(0),Inches(0),P.slide_width,Inches(1.15),D);Tx(s,Inches(0.8),Inches(0.2),Inches(11),Inches(0.7),tx,fs=32,c=W,b=True);Bx(s,Inches(0.8),Inches(1.0),Inches(2.5),Inches(0.05),O)

# ====== S1: Title ======
s=P.slides.add_slide(P.slide_layouts[6])
s.background.fill.solid();s.background.fill.fore_color.rgb=D
Bx(s,Inches(0),Inches(0),P.slide_width,Inches(0.12),O)
Bx(s,Inches(0),Inches(7.38),P.slide_width,Inches(0.12),O)
Tx(s,Inches(1),Inches(1.5),Inches(11),Inches(1.2),'校园RPG冒险游戏',fs=52,c=W,b=True,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(2.8),Inches(11),Inches(0.7),'C++面向对象程序设计 · 课程设计汇报',fs=26,c=O,a=PP_ALIGN.CENTER)
Bx(s,Inches(5),Inches(3.8),Inches(3.3),Inches(0.03),O)
Tx(s,Inches(1),Inches(4.2),Inches(11),Inches(0.5),'C++17 + Win32 API + MinGW GCC 6.3.0',fs=18,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(5.4),Inches(11),Inches(0.5),'组  长：张佳怡    学  号：1030425404',fs=20,c=W,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(6.0),Inches(11),Inches(0.5),'人工智能与自动化学院 · 计算机科学与技术 · 2026年7月',fs=16,c=Y,a=PP_ALIGN.CENTER)

# ====== S2: TOC ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'汇报提纲')
toc=['01  项目背景与目标','02  项目特色与创新（2项挑战任务）','03  项目管理（团队分工+进度安排）',
     '04  需求分析（功能+非功能需求）','05  系统设计（架构+模块+数据流）','06  类设计与UML（继承+组合）',
     '07  关键功能实现（代码详解）','08  测试与结果分析','09  总结与展望']
for i,item in enumerate(toc):
 y=Inches(1.6)+Inches(0.6)*i
 Bx(s,Inches(1.2),y,Inches(0.55),Inches(0.45),A if i<2 else L)
 Tx(s,Inches(1.2),y+Inches(0.05),Inches(0.55),Inches(0.35),item[:2],fs=16,c=W,b=True,a=PP_ALIGN.CENTER)
 Tx(s,Inches(2.0),y+Inches(0.05),Inches(9),Inches(0.35),item[3:],fs=19,c=K)

# ====== S3: Project Overview ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目背景与目标')
Li(s,Inches(0.8),Inches(1.6),Inches(7.5),Inches(5.5),[
'• 项目类型：基于C++17的Windows GUI校园冒险RPG游戏',
'• 游戏世界观：以校园生活为冒险背景，轻松休闲的角色扮演世界',
'• 玩家角色：扮演校园冒险者，探索校园、遭遇怪物、完成任务',
'• 开发模式：小组合作（4人），面向对象分析设计与编程',
'• 技术栈：C++17 + Win32 API + MinGW GCC 6.3.0 + CMake',
'• 核心特性：类继承、多态、虚函数、STL容器、文件持久化、多线程',
'• 代码规模：约3000行源代码，包含15个核心类，6大功能模块',
'• 界面风格：矩形按钮交互，微软雅黑字体，中文界面',
'• 编译方案：-fexec-charset=UTF-8 + CP_UTF8编码页解决中文乱码',
],fs=15,sp=Pt(10))
nums=[('15+','核心类'),('3000+','代码行数'),('6','功能模块'),('12','游戏道具'),('4','校园任务'),('3','怪物类型')]
for i,(n,label) in enumerate(nums):
 y=Inches(1.6)+Inches(0.95)*i
 Bx(s,Inches(8.8),y,Inches(4.0),Inches(0.8),G)
 Tx(s,Inches(9.0),y+Inches(0.05),Inches(1.5),Inches(0.6),n,fs=32,c=A,b=True)
 Tx(s,Inches(10.5),y+Inches(0.2),Inches(2.0),Inches(0.4),label,fs=16,c=K)

# ====== S4: Innovation ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目特色与创新 — 2项挑战任务')
# Card 1
Bx(s,Inches(0.8),Inches(1.8),Inches(5.8),Inches(2.3),G)
Tx(s,Inches(1.1),Inches(1.9),Inches(5.3),Inches(0.4),'挑战1：Item多态继承体系',fs=22,c=D,b=True)
Bx(s,Inches(1.1),Inches(2.35),Inches(5.3),Inches(0.03),A)
Li(s,Inches(1.1),Inches(2.5),Inches(5.3),Inches(1.5),[
'• Item抽象基类：纯虚函数 use()/clone()/getTypeName()',
'• Food(食物)：恢复HP 25~65，3种校园食品',
'• Medicine(药品)：高额回血 50~9999，3种医疗用品',
'• Equipment(装备)：永久ATK/DEF加成，武器/防具/饰品',
'• 背包统一用 vector<Item*> 基类指针管理，运行时多态分派',
],fs=13,sp=Pt(5))
# Card 2
Bx(s,Inches(7.0),Inches(1.8),Inches(5.5),Inches(2.3),G)
Tx(s,Inches(7.3),Inches(1.9),Inches(5.0),Inches(0.4),'挑战2：Skill技能系统',fs=22,c=D,b=True)
Bx(s,Inches(7.3),Inches(2.35),Inches(5.0),Inches(0.03),A)
Li(s,Inches(7.3),Inches(2.5),Inches(5.0),Inches(1.5),[
'• Skill抽象基类：纯虚函数 calcDamage()/clone()',
'• MeleeSkill(物理)：伤害=力量×系数-目标物防',
'• MagicSkill(魔法)：伤害=智力×系数-目标魔防',
'• HealSkill(治疗)：恢复HP=智力×N (4/5/6)',
'• DebuffSkill(减益)：伤害+降低敌方物防30%',
],fs=13,sp=Pt(5))
# Bottom
Tx(s,Inches(0.8),Inches(4.5),Inches(11.5),Inches(0.4),'设计模式应用',fs=18,c=D,b=True)
patterns=['工厂模式(SkillFactory)  |  原型模式(Item::clone())  |  策略模式(Skill::calcDamage)  |  模板方法(序列化接口)']
Li(s,Inches(0.8),Inches(5.0),Inches(11.5),Inches(1.5),patterns,fs=15,c=A)

# ====== S5: Challenge 1 Detail ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'挑战1详解 — Item多态继承体系')
# Class tree
Tx(s,Inches(0.8),Inches(1.5),Inches(6),Inches(0.4),'类层次结构 & 12种道具',fs=20,c=D,b=True)
Bx(s,Inches(0.8),Inches(2.0),Inches(6),Inches(4.8),G)
classes=[
 ('Item (抽象基类)','id/name/desc/type/price/quantity | 纯虚: use()/clone()/getTypeName()',D),
 ('Food : Item','m_hpRestore | use(): target.addHP(25~65) | getTypeName(): 食物',A),
 ('Medicine : Item','m_hpRestore + m_cureEffect | use(): target.addHP(50~9999) | getTypeName(): 药品',A),
 ('Equipment : Item','m_attackBonus/m_defenseBonus/m_slot | use(): 永久属性加成(武器/防具/饰品)',A),
]
for i,(name,detail,clr) in enumerate(classes):
 y=Inches(2.15)+Inches(1.15)*i
 Bx(s,Inches(1.0),y,Inches(5.6),Inches(1.05),W)
 Tx(s,Inches(1.15),y+Inches(0.03),Inches(5.3),Inches(0.3),name,fs=14,c=clr,b=True)
 Tx(s,Inches(1.15),y+Inches(0.33),Inches(5.3),Inches(0.65),detail,fs=11,c=K)
# Right: items table
Tx(s,Inches(7.3),Inches(1.5),Inches(5.5),Inches(0.4),'12种校园道具（全部可在商店购买）',fs=18,c=D,b=True)
items_list=[
 '【食物类】校园面包(15G,HP+25) | 清甜牛奶(25G,HP+40) | 能量饭团(40G,HP+65)',
 '【药品类】基础创可贴(30G,HP+50) | 急救喷雾(60G,HP+100) | 全能药剂(120G,满血)',
 '【装备类】新手校服(80G,DEF+8) | 运动护腕(150G,ATK+12) | 学霸徽章(280G,ATK+15/DEF+10)',
]
for i,item in enumerate(items_list):
 Bx(s,Inches(7.3),Inches(2.1)+Inches(1.2)*i,Inches(5.5),Inches(1.0),G)
 Tx(s,Inches(7.5),Inches(2.2)+Inches(1.2)*i,Inches(5.0),Inches(0.8),item,fs=13,c=K)
# Code
Tx(s,Inches(7.3),Inches(5.8),Inches(5.5),Inches(0.3),'核心代码',fs=14,c=D,b=True)
code_block=Bx(s,Inches(7.3),Inches(6.1),Inches(5.5),Inches(1.1),RGBColor(0x1E,0x1E,0x1E))
Tx(s,Inches(7.5),Inches(6.15),Inches(5.0),Inches(1.0),'class Item {\n  virtual void use(BaseCharacter&)=0;\n  virtual Item* clone() const=0;\n  virtual string getTypeName() const=0;\n};',fs=10,c=RGBColor(0xD4,0xD4,0xD4))

# ====== S6: Challenge 2 Detail ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'挑战2详解 — Skill技能系统')
Tx(s,Inches(0.8),Inches(1.5),Inches(6),Inches(0.4),'技能类层次结构',fs=20,c=D,b=True)
skills=[
 ('Skill (抽象基类)','id/name/lv1~3/mpCost/coeff | 纯虚: calcDamage()/clone()',D),
 ('MeleeSkill : Skill','物理伤害=力量×系数-物防 | Lv1:×1.2 Lv2:×1.4 Lv3:×1.7',A),
 ('MagicSkill : Skill','魔法伤害=智力×系数-魔防 | Lv1:×1.3 Lv2:×1.5 Lv3:×1.8',A),
 ('HealSkill : Skill','恢复HP=智力×N | Lv1:N=4 Lv2:N=5 Lv3:N=6',A),
 ('DebuffSkill : Skill','0.8倍魔法伤害 + 降低敌方物防30% 持续1回合',A),
]
for i,(name,detail,clr) in enumerate(skills):
 y=Inches(2.0)+Inches(1.0)*i
 Bx(s,Inches(0.8),y,Inches(6.2),Inches(0.9),G if i>0 else W)
 Tx(s,Inches(1.0),y+Inches(0.03),Inches(5.8),Inches(0.28),name,fs=14,c=clr,b=True)
 Tx(s,Inches(1.0),y+Inches(0.3),Inches(5.8),Inches(0.55),detail,fs=11,c=K)
# Right: 6 skills
Tx(s,Inches(7.5),Inches(1.5),Inches(5.3),Inches(0.4),'6种校园主题技能',fs=20,c=D,b=True)
skill_data=[
 ('1. 挥拳重击','Melee | MP:10 | 解锁Lv1','基础物理攻击，1.2~1.7倍力量伤害'),
 ('2. 冲刺冲撞','Melee | MP:22 | 解锁Lv3','强力冲撞，1.8~2.5倍高额物理伤害'),
 ('3. 粉笔飞弹','Magic | MP:12 | 解锁Lv1','基础魔法攻击，1.3~1.8倍智力伤害'),
 ('4. 试卷风暴','Magic | MP:30 | 解锁Lv4','大范围魔法攻击，2.2~3.0倍智力伤害'),
 ('5. 课间补给','Heal | MP:15 | 解锁Lv2','恢复HP = 智力×4~6，战斗续航'),
 ('6. 说教威慑','Debuff | MP:18 | 解锁Lv3','0.8倍魔法伤害 + 降低敌方物防30%'),
]
for i,(name,meta,desc) in enumerate(skill_data):
 y=Inches(2.0)+Inches(0.82)*i
 Bx(s,Inches(7.5),y,Inches(5.3),Inches(0.72),G)
 Tx(s,Inches(7.7),y+Inches(0.02),Inches(2.5),Inches(0.3),name,fs=13,c=A,b=True)
 Tx(s,Inches(10.0),y+Inches(0.02),Inches(2.6),Inches(0.3),meta,fs=10,c=K)
 Tx(s,Inches(7.7),y+Inches(0.35),Inches(5.0),Inches(0.3),desc,fs=10,c=K)
Tx(s,Inches(7.5),Inches(6.95),Inches(5.3),Inches(0.3),'SkillFactory::createSkill(id) → clone() 工厂模式',fs=12,c=A,b=True)

# ====== S7: Team ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目管理 — 团队分工与进度安排')
# Table
rows,cols=5,6
tbl=s.shapes.add_table(rows,cols,Inches(0.5),Inches(1.6),Inches(12.3),Inches(2.4)).table
hdrs=['序号','角色','姓名','学号','班级','职责分工']
data=[
 ['1','组长','张佳怡','1030425404','计科2201','项目架构、Item多态体系(基类+3派生类)、背包系统、商店系统、CMake构建、编码修复'],
 ['2','组员','成员B','xxxxxxxxx','计科2201','Win32 GUI框架、resources.rc、回合制战斗系统(3种怪物)、怪物属性设计'],
 ['3','组员','成员C','xxxxxxxxx','计科2201','角色系统(BaseCharacter+Player)、等级成长、存档系统、多线程自动存档'],
 ['4','组员','成员D','xxxxxxxxx','计科2201','任务系统(Task+TaskSystem)、任务模板文件、技能系统、SkillFactory、文档撰写'],
]
for j,h in enumerate(hdrs):
 c=tbl.cell(0,j);c.text=h
 for p in c.text_frame.paragraphs:p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=W;p.alignment=PP_ALIGN.CENTER
 c.fill.solid();c.fill.fore_color.rgb=D
for i,rd in enumerate(data):
 for j,tx in enumerate(rd):
  c=tbl.cell(i+1,j);c.text=tx
  for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.alignment=PP_ALIGN.CENTER
  if i%2==0:c.fill.solid();c.fill.fore_color.rgb=G
# Timeline
Tx(s,Inches(0.8),Inches(4.3),Inches(5),Inches(0.4),'项目进度甘特图（12周）',fs=18,c=D,b=True)
phases=['需求分析\nUML建模','核心类实现\nItem/Skill','进阶功能\n战斗/任务','GUI开发\nWin32','存档集成\n多线程','测试与\n文档撰写']
weeks=['第1-2周','第3-4周','第5-6周','第7-8周','第9-10周','第11-12周']
for i in range(6):
 x=Inches(0.5)+Inches(2.1)*i
 Bx(s,x,Inches(4.8),Inches(1.9),Inches(1.1),A if i<2 else L)
 Tx(s,x,Inches(4.85),Inches(1.9),Inches(0.5),phases[i],fs=11,c=W,b=True,a=PP_ALIGN.CENTER)
 Tx(s,x,Inches(5.4),Inches(1.9),Inches(0.3),weeks[i],fs=9,c=W,a=PP_ALIGN.CENTER)
 if i<5:Bx(s,x+Inches(1.8),Inches(5.3),Inches(0.35),Inches(0.06),D)
Tx(s,Inches(0.8),Inches(6.3),Inches(11.5),Inches(0.8),'协作方式：Git版本控制 + 模块化开发 + 接口约定(头文件先行) + 定期Code Review',fs=14,c=A,b=True)

# ====== S8: Requirements ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'需求分析 — 功能需求')
# 6 function cards
funcs=[
 ('角色管理','创建/查看/存读档\n名称、等级Lv1~50\nHP/MP/EXP/金币\nSTR/VIT/AGI/INT',A),
 ('背包管理','获得/查看/使用/删除\n3类12种道具\n食物/药品/装备\n无限容量',A),
 ('商店系统','查看/购买/出售\n12种固定商品\n半价回收结算\n实时金币更新',A),
 ('任务系统','查看/接受/完成/领奖\n4个梯度任务\n击杀计数追踪\n状态机流转',D),
 ('战斗系统','选择/攻击/受击\n3种校园怪物\n回合制自动战斗\nEXP+金币奖励',D),
 ('等级成长','经验累计\n自动检测升级\nHP+20/ATK+5每级\n最高50级上限',D),
]
for i,(title,desc,clr) in enumerate(funcs):
 x=Inches(0.5)+Inches(2.15)*i
 Bx(s,x,Inches(1.6),Inches(2.0),Inches(3.2),G)
 Tx(s,x+Inches(0.1),Inches(1.7),Inches(1.8),Inches(0.4),title,fs=16,c=clr,b=True,a=PP_ALIGN.CENTER)
 Bx(s,x+Inches(0.15),Inches(2.15),Inches(1.7),Inches(0.03),clr)
 Tx(s,x+Inches(0.1),Inches(2.3),Inches(1.8),Inches(2.3),desc,fs=13,c=K,a=PP_ALIGN.CENTER)
# Non-functional
Tx(s,Inches(0.8),Inches(5.1),Inches(11.5),Inches(0.3),'非功能需求',fs=18,c=D,b=True)
Li(s,Inches(0.8),Inches(5.4),Inches(11.5),Inches(1.5),[
'• 性能：GUI响应<500ms | 战斗结算<200ms | 存档写入<100ms',
'• 可靠性：多线程自动存档每60秒 | volatile标志位线程安全 | 存档文件校验',
'• 可用性：微软雅黑字体渲染 | 中文界面无乱码 | 矩形按钮直观交互',
'• 可维护性：头文件与实现分离 | 命名空间RPG | 清晰的注释规范',
'• 可扩展性：虚函数接口继承扩展 | 工厂模式注册新类型 | CMake模块化构建',
'• 兼容性：Windows 7/10/11 | MinGW GCC 6.3.0+ | 静态链接无额外依赖',
],fs=12,sp=Pt(4))

# ====== S9: Architecture ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'系统设计 — 三层架构')
layers=[
 ('表示层 — Win32 GUI',A,'main.cpp (619行)  WndProc消息循环  9按钮事件分发\n8个模态对话框(DialogBoxW)  resources.rc资源文件\n微软雅黑GDI字体渲染  状态栏实时信息展示'),
 ('—— 全局变量交互 g_player/g_inventory/g_shop/g_taskSystem ——','',''),
 ('业务逻辑层 — 15个核心类',D,'Player角色属性/等级/技能/金币  Inventory背包(vector<Item*>)\nShop商店(12商品模板)  TaskSystem任务管理(状态机)\nSkill体系(6技能工厂)  Item体系(12道具多态)'),
 ('—— 序列化接口 serialize/deserialize ——','',''),
 ('数据持久层 — 文件I/O + 多线程',L,'SaveManager玩家数据存读档  tasks.txt任务模板(UTF-8)\nWin32 CreateThread后台线程  每60秒自动保存\nvolatile标志位线程通信  WaitForSingleObject安全退出'),
]
for i,(name,clr,desc) in enumerate(layers):
 if '——' in name:
  Tx(s,Inches(1.2),Inches(1.2)+Inches(1.0)*i,Inches(10.8),Inches(0.35),name[:50],fs=11,c=O,b=True,a=PP_ALIGN.CENTER)
  continue
 y=Inches(1.3)+Inches(1.1)*i
 Bx(s,Inches(1.0),y,Inches(11.3),Inches(1.3),G)
 Tx(s,Inches(1.3),y+Inches(0.05),Inches(10.7),Inches(0.35),name,fs=17,c=clr,b=True)
 Tx(s,Inches(1.3),y+Inches(0.42),Inches(10.7),Inches(0.8),desc,fs=12,c=K)

# ====== S10: Module Design ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'系统设计 — 六大模块详解')
mods=[
 ('角色模块','BaseCharacter(抽象基类)\n• 通用属性/战斗接口\n• 虚函数serialize()\n\nPlayer(派生类)\n• exp/gold/skillPoint\n• 技能/背包/任务容器\n• 等级成长/金币操作',A),
 ('道具模块','Item(抽象基类)\n• use()/clone()纯虚函数\n\nFood(食物)\n• HP恢复25~65\n\nMedicine(药品)\n• HP恢复50~9999\n\nEquipment(装备)\n• 永久ATK/DEF加成',D),
 ('背包+商店','Inventory(背包)\n• vector<Item*>管理\n• 增删查改/序列化\n\nShop(商店)\n• 12种商品模板\n• 购买clone/出售半价\n• 自动金币结算',L),
 ('任务模块','TaskSystem\n• 4个任务模板\n• 从tasks.txt加载\n\nTask(状态机)\nNotAccepted→InProgress\n→Completed→RewardClaimed\n• 击杀计数追踪',A),
 ('战斗模块','EnemyInfo结构体\n3种校园怪物：\n• 懒散学渣(Lv1)\n• 熬夜复习怪(Lv2)\n• 期末Boss(Lv3)\n\nwhile回合制循环\nrand()随机伤害浮动',D),
 ('技能模块','SkillFactory(工厂)\n管理6种技能模板\nclone()创建实例\n\n4类派生Skill：\nMelee/Magic/Heal/Debuff\ncalcDamage()策略模式',L),
]
for i,(title,detail,clr) in enumerate(mods):
 x=Inches(0.3)+Inches(2.2)*i
 Bx(s,x,Inches(1.5),Inches(2.05),Inches(5.5),G)
 Tx(s,x+Inches(0.1),Inches(1.6),Inches(1.85),Inches(0.35),title,fs=16,c=clr,b=True,a=PP_ALIGN.CENTER)
 Bx(s,x+Inches(0.15),Inches(2.0),Inches(1.75),Inches(0.03),O)
 Tx(s,x+Inches(0.1),Inches(2.15),Inches(1.85),Inches(4.5),detail,fs=11,c=K)

# ====== S11: UML ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'类设计与UML — 完整类图')
Tx(s,Inches(0.8),Inches(1.5),Inches(5.8),Inches(0.4),'继承关系 Inheritance',fs=20,c=D,b=True)
Bx(s,Inches(0.8),Inches(2.0),Inches(5.8),Inches(2.5),G)
Tx(s,Inches(1.0),Inches(2.1),Inches(5.3),Inches(0.25),'Item多态体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(2.4),Inches(5.3),Inches(0.22),'Item ──→ Food ──→ Medicine ──→ Equipment',fs=12,c=K)
Tx(s,Inches(1.0),Inches(2.6),Inches(5.3),Inches(0.2),'(抽象基类)    (HP恢复)       (高额回血)        (永久加成)',fs=9,c=K)
Tx(s,Inches(1.0),Inches(2.9),Inches(5.3),Inches(0.25),'Skill多态体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(3.2),Inches(5.3),Inches(0.22),'Skill ──→ MeleeSkill / MagicSkill / HealSkill / DebuffSkill',fs=12,c=K)
Tx(s,Inches(1.0),Inches(3.4),Inches(5.3),Inches(0.2),'(抽象基类)       (物理)          (魔法)        (治疗)         (减益)',fs=9,c=K)
Tx(s,Inches(1.0),Inches(3.7),Inches(5.3),Inches(0.25),'角色类体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(4.0),Inches(5.3),Inches(0.22),'BaseCharacter ──→ Player (技能/背包/任务容器)',fs=12,c=K)
# Right: composition
Tx(s,Inches(7.0),Inches(1.5),Inches(5.5),Inches(0.4),'组合/聚合关系 Composition',fs=20,c=D,b=True)
Bx(s,Inches(7.0),Inches(2.0),Inches(5.5),Inches(2.5),G)
Li(s,Inches(7.2),Inches(2.1),Inches(5.0),Inches(2.3),[
'• Inventory ◇──→ Item*  (vector容器)',
'• Shop     ◇──→ Item*  (12种商品模板)',
'• TaskSystem ◇──→ Task  (vector<Task>)',
'• Player   ◇──→ Skill* (已学技能列表)',
'• SaveManager ..→ Player (序列化依赖)',
'• SkillFactory ..→ Skill (工厂创建)',
],fs=14,c=K,sp=Pt(12))
# Summary
Tx(s,Inches(0.8),Inches(4.8),Inches(11.5),Inches(0.4),'项目统计：15个核心类 | 3个抽象基类 | 4层最大继承深度 | 5种关联关系',fs=16,c=A,b=True,a=PP_ALIGN.CENTER)
# Code diagram
Bx(s,Inches(0.8),Inches(5.3),Inches(11.7),Inches(1.9),RGBColor(0x1E,0x1E,0x1E))
Tx(s,Inches(1.0),Inches(5.4),Inches(11.3),Inches(1.7),
   '  BaseCharacter <|--- Player\n'
   '       ↑                           Inventory ◇--- Item*\n'
   '  Item <|--- Food <|--- Medicine <|--- Equipment    Shop ◇--- Item*\n'
   '  Skill <|--- MeleeSkill / MagicSkill               TaskSystem ◇--- Task\n'
   '       <|--- HealSkill / DebuffSkill                Player ◇--- Skill*',
   fs=11,c=RGBColor(0x4E,0xC9,0xB0))

# ====== S12: Key Implementation 1 ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'关键功能实现 (1/2)')
impls=[
 ('1. 物品多态 — use()运行时分派',
  'Item基类定义 pure virtual use(BaseCharacter&)=0。三个派生类各实现不同use()策略：\n'
  'Food→target.addHP(m_hpRestore) | Medicine→target.addHP(高额) | Equipment→永久addAttack/addDefense\n'
  '背包使用vector<Item*>统一存储，遍历调用it->use(player)时C++虚函数表自动分派到正确实现。\n'
  '商店购买通过clone()原型模式创建物品副本，避免直接操作商店模板。'),
 ('2. 回合制战斗 — while循环模拟',
  'EnemyInfo结构体存储怪物(name/hp/atk/def/exp/gold)。while(pHp>0&&eHp>0&&rounds<50)：\n'
  '玩家伤害=ATK+rand()%(ATK/2+1)-eDef | 敌方伤害=eAtk+rand()%(eAtk/2+1)-pDef | 最小值=1\n'
  '50回合上限防死循环。胜利→addExp+gainGold+击杀计数++→checkLevelUp()。失败→HP=1，无惩罚。'),
 ('3. 多线程自动存档 — Win32 CreateThread',
  'MinGW 6.3.0不支持std::thread → CreateThread(NULL,0,proc,NULL,0,NULL)创建后台线程。\n'
  'autoSaveLoop(): for(int i=0;i<60;i++)Sleep(1000); if(g_autoSaveRunning)SaveManager::save()\n'
  'volatile bool控制线程生命周期。退出: g_autoSaveRunning=false→WaitForSingleObject(5000)→CloseHandle。'),
]
for i,(title,desc) in enumerate(impls):
 y=Inches(1.4)+Inches(2.0)*i
 Bx(s,Inches(0.5),y,Inches(0.08),Inches(1.7),O)
 Tx(s,Inches(0.9),y-Inches(0.05),Inches(11.5),Inches(0.4),title,fs=17,c=D,b=True)
 Tx(s,Inches(0.9),y+Inches(0.4),Inches(11.8),Inches(1.4),desc,fs=13,c=K)

# ====== S13: Key Implementation 2 ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'关键功能实现 (2/2)')
impls2=[
 ('4. 等级自动成长 — checkLevelUp()',
  '每次获得经验后调用checkLevelUp()。while循环：\n'
  'if(getExp()>=expToNextLevel()){exp-=need;level++;maxHp+=20;atk+=5;hp=maxHp;recalcStats();}\n'
  '阈值公式：expToNextLevel(lv)=lv×30+20。Lv1→2需50EXP，Lv49→50需1490EXP。\n'
  '最高50级。溢出经验自动保留。升级时满血满蓝。全程自动化无需玩家手动操作。'),
 ('5. UTF-8编码转换 — 解决中文乱码',
  '编译选项-fexec-charset=UTF-8使窄字符串以UTF-8存入二进制。Win32 W-API需UTF-16 LE。\n'
  '原U2W/W2U/SetStatus使用CP_ACP(GBK/CP936)→中文乱码。\n'
  '解决方案：MultiByteToWideChar/WideCharToMultiByte的CodePage参数改为CP_UTF8。\n'
  '配合-finput-charset=UTF-8源文件编码，实现完整的UTF-8↔UTF-16正确互转。'),
 ('6. Win32 GUI架构 — 消息驱动模型',
  'WinMain: 注册窗口类→CreateWindowExW创建主窗口→GetMessage/DispatchMessage消息循环。\n'
  'WndProc: WM_CREATE(9按钮+状态栏+字体) | WM_COMMAND(按钮→DialogBoxW模态对话框)\n'
  'WM_CLOSE→DestroyWindow | WM_DESTROY→清理资源+PostQuitMessage。\n'
  'resources.rc(GBK编码)定义8个对话框布局，windres编译为.o文件链接。'),
]
for i,(title,desc) in enumerate(impls2):
 y=Inches(1.4)+Inches(2.0)*i
 Bx(s,Inches(0.5),y,Inches(0.08),Inches(1.7),O)
 Tx(s,Inches(0.9),y-Inches(0.05),Inches(11.5),Inches(0.4),title,fs=17,c=D,b=True)
 Tx(s,Inches(0.9),y+Inches(0.4),Inches(11.8),Inches(1.4),desc,fs=13,c=K)

# ====== S14: Combat Detail ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'战斗系统 — 详细流程与怪物属性')
# Flow
Tx(s,Inches(0.8),Inches(1.5),Inches(5.5),Inches(0.4),'战斗流程',fs=18,c=D,b=True)
steps=[
 '1. 玩家点击[挑战战斗]→弹出对话框',
 '2. 选择敌人(懒散学渣/熬夜怪/期末Boss)',
 '3. 检查玩家HP>0（否则提示恢复HP）',
 '4. while(pHp>0 && eHp>0 && rounds<50){',
 '5.   玩家伤害=ATK+rand()%(ATK/2+1)-eDef',
 '6.   敌方HP-=max(1,玩家伤害)',
 '7.   if(敌方HP<=0) break; // 胜利!',
 '8.   敌方伤害=eAtk+rand()%(eAtk/2+1)-pDef',
 '9.   玩家HP-=max(1,敌方伤害)',
 '10.}',
 '11. 胜利→addExp+gainGold+击杀计数+checkLevelUp',
 '12. 失败→HP设为1，无惩罚结束战斗',
]
Li(s,Inches(0.8),Inches(2.0),Inches(5.5),Inches(5.0),steps,fs=11,c=K,sp=Pt(3))
# Enemy stats
Tx(s,Inches(7.0),Inches(1.5),Inches(5.8),Inches(0.4),'三种校园怪物属性',fs=18,c=D,b=True)
enemies=[
 ('懒散学渣小怪 (新手级)','HP:80 | ATK:8 | DEF:2','EXP:30 | 金币:15','校园最弱怪物，适合新手练级'),
 ('熬夜复习怪 (中等级)','HP:150 | ATK:15 | DEF:5','EXP:65 | 金币:35','长期熬夜积攒的负能量怪物'),
 ('期末压力Boss (高级)','HP:300 | ATK:28 | DEF:10','EXP:150 | 金币:90','校园最强BOSS，奖励最丰厚'),
]
for i,(name,stats,reward,desc) in enumerate(enemies):
 y=Inches(2.0)+Inches(1.7)*i
 Bx(s,Inches(7.0),y,Inches(5.8),Inches(1.5),G)
 Tx(s,Inches(7.2),y+Inches(0.05),Inches(5.3),Inches(0.3),name,fs=15,c=A,b=True)
 Tx(s,Inches(7.2),y+Inches(0.4),Inches(2.5),Inches(0.25),stats,fs=12,c=K)
 Tx(s,Inches(9.8),y+Inches(0.4),Inches(2.8),Inches(0.25),reward,fs=12,c=O)
 Tx(s,Inches(7.2),y+Inches(0.75),Inches(5.3),Inches(0.25),desc,fs=11,c=K)

# ====== S15: Task System ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'任务系统 — 4个校园任务')
# State machine
Tx(s,Inches(0.8),Inches(1.5),Inches(5.5),Inches(0.4),'任务状态机',fs=18,c=D,b=True)
states=['NotAccepted\n(未接受)','→','InProgress\n(进行中)','→','Completed\n(已完成)','→','RewardClaimed\n(已领奖)']
for i,st in enumerate(states):
 x=Inches(0.3)+Inches(1.9)*i
 if st=='→':
  Tx(s,x,Inches(2.2),Inches(0.5),Inches(0.4),'→',fs=24,c=O,b=True,a=PP_ALIGN.CENTER)
 else:
  clr=D if i==0 else (A if i==2 else (L if i==4 else RGBColor(0x88,0x88,0x88)))
  Bx(s,x,Inches(2.0),Inches(1.7),Inches(0.8),clr)
  Tx(s,x,Inches(2.1),Inches(1.7),Inches(0.6),st,fs=10,c=W,b=True,a=PP_ALIGN.CENTER)
# 4 tasks
Tx(s,Inches(0.8),Inches(3.0),Inches(6.0),Inches(0.4),'任务详情',fs=18,c=D,b=True)
tasks=[
 ('任务一：新手初试冒险','完成1场任意战斗','EXP:80 | 金币:50 | 校园面包×2',O),
 ('任务二：清理校园小怪','累计击杀3只懒散学渣','EXP:120 | 金币:80 | 急救喷雾×1',O),
 ('任务三：中级历练挑战','击败2只熬夜复习怪','EXP:200 | 金币:120 | 运动护腕×1',O),
 ('任务四：攻克期末难关','击败1只期末压力Boss','EXP:350 | 金币:200 | 全能药剂×2+学霸徽章',O),
]
for i,(name,cond,reward,clr) in enumerate(tasks):
 y=Inches(3.5)+Inches(0.85)*i
 Bx(s,Inches(0.8),y,Inches(6.0),Inches(0.7),G)
 Tx(s,Inches(1.0),y+Inches(0.05),Inches(2.8),Inches(0.25),name,fs=14,c=clr,b=True)
 Tx(s,Inches(3.9),y+Inches(0.05),Inches(2.7),Inches(0.25),cond,fs=11,c=K)
 Tx(s,Inches(1.0),y+Inches(0.35),Inches(5.5),Inches(0.25),reward,fs=10,c=A)
# Right: kill tracking
Tx(s,Inches(7.3),Inches(1.5),Inches(5.5),Inches(0.4),'击杀计数联动机制',fs=18,c=D,b=True)
Bx(s,Inches(7.3),Inches(2.0),Inches(5.5),Inches(2.5),RGBColor(0x1E,0x1E,0x1E))
Tx(s,Inches(7.5),Inches(2.05),Inches(5.0),Inches(2.3),
   'int g_killSlacker=0;   // 懒散学渣\n'
   'int g_killNightOwl=0;  // 熬夜复习怪\n'
   'int g_killBoss=0;      // 期末Boss\n'
   'int g_totalBattles=0;  // 总战斗数\n\n'
   '// 战斗胜利后自动更新\n'
   'if(sel==0) g_killSlacker++;\n'
   'if(sel==1) g_killNightOwl++;\n'
   'if(sel==2) g_killBoss++;\n\n'
   '// 任务完成判定\n'
   'case0: canComplete=(totalBattles>=1)\n'
   'case1: canComplete=(killSlacker>=3)\n'
   'case2: canComplete=(killNightOwl>=2)\n'
   'case3: canComplete=(killBoss>=1)',
   fs=10,c=RGBColor(0x4E,0xC9,0xB0))
Tx(s,Inches(7.3),Inches(4.7),Inches(5.5),Inches(0.3),'tasks.txt 文件格式 (UTF-8编码，管道符分隔)',fs=14,c=D,b=True)
Tx(s,Inches(7.3),Inches(5.1),Inches(5.5),Inches(1.5),
   '名称|描述|条件|经验|金币\n'
   '新手初试冒险|初次开启校园冒险...|完成1场战斗|80|50\n'
   '清理校园小怪|维护校园和平...|击杀3只懒散学渣|120|80\n'
   '中级历练挑战|挑战中级怪物...|击败2只熬夜复习怪|200|120\n'
   '攻克期末难关|挑战校园最强BOSS...|击败1只期末Boss|350|200',
   fs=11,c=K)

# ====== S16: Testing ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'测试与结果分析')
tests=[
 ('角色管理','创建角色输入中文名','名称正确显示','✓'),
 ('角色管理','查看属性面板','Lv.1数据全部正确','✓'),
 ('背包管理','打开空背包','显示"背包是空的"','✓'),
 ('商店系统','购买校园面包(15G)','金币-15,背包+1','✓'),
 ('商店系统','出售物品(半价回收)','金币+7,物品移除','✓'),
 ('商店系统','金币不足时购买','提示"金币不足"','✓'),
 ('战斗系统','挑战懒散学渣小怪','回合制+30EXP+15G','✓'),
 ('战斗系统','挑战期末压力Boss','150EXP+90G奖励','✓'),
 ('战斗系统','玩家HP归零战斗失败','HP恢复为1，无惩罚','✓'),
 ('等级成长','EXP累积达到阈值','自动升级,HP+20,ATK+5','✓'),
 ('任务系统','接取→完成→领取奖励','状态流转+奖励发放正确','✓'),
 ('存档系统','保存后关闭重新打开','等级/金币/背包数据恢复','✓'),
 ('中文显示','全界面物品/怪物/任务名','无任何中文乱码','✓'),
]
for i,(mod,case_text,exp_text,res) in enumerate(tests):
 y=Inches(1.4)+Inches(0.42)*i
 if i%2==0:Bx(s,Inches(0.5),y,Inches(8.2),Inches(0.38),G)
 Tx(s,Inches(0.6),y+Inches(0.02),Inches(1.6),Inches(0.32),mod,fs=11,c=A,b=True)
 Tx(s,Inches(2.3),y+Inches(0.02),Inches(3.5),Inches(0.32),case_text,fs=11,c=K)
 Tx(s,Inches(5.9),y+Inches(0.02),Inches(2.0),Inches(0.32),exp_text,fs=10,c=K)
 Bx(s,Inches(8.1),y+Inches(0.03),Inches(0.5),Inches(0.3),R)
 Tx(s,Inches(8.1),y+Inches(0.03),Inches(0.5),Inches(0.3),res,fs=14,c=W,b=True,a=PP_ALIGN.CENTER)
# Right: Problems
Tx(s,Inches(9.0),Inches(1.5),Inches(4.0),Inches(0.4),'关键问题与解决方案',fs=16,c=D,b=True)
probs=[
 ('问题1：中文乱码','编译选项-fexec-charset=UTF-8输出UTF-8窄字符串，但U2W/W2U使用CP_ACP(GBK)导致乱码。\n→ 将MultiByteToWideChar的CodePage从CP_ACP改为CP_UTF8'),
 ('问题2：调试对话框bug','商店初始化代码中含wsprintfW(%s传入char*)类型不匹配的调试消息框。\n→ 移除调试代码'),
 ('问题3：构建工具缺失','MinGW安装不含mingw32-make，CMake无法生成Makefile。\n→ 改用g++命令行直接编译，更新CMakeLists.txt补充所有源文件'),
 ('问题4：编码策略统一','L"..."宽字符串可正常显示但窄字符串需U2W转换。\n→ 确认统一使用UTF-8编码策略(-fexec-charset=UTF-8 + CP_UTF8)'),
]
for i,(prob,solve) in enumerate(probs):
 y=Inches(2.1)+Inches(1.25)*i
 Bx(s,Inches(9.0),y,Inches(4.0),Inches(1.1),G)
 Tx(s,Inches(9.2),y+Inches(0.05),Inches(3.6),Inches(0.3),prob,fs=12,c=RGBColor(0xC0,0x39,0x2B),b=True)
 Tx(s,Inches(9.2),y+Inches(0.35),Inches(3.6),Inches(0.65),solve,fs=10,c=K)
# Bottom
Bx(s,Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.4),A)
Tx(s,Inches(0.8),Inches(6.87),Inches(11.7),Inches(0.35),'测试统计：13个测试用例 | 13个通过 | 0个失败 | 通过率 100% | 覆盖6大模块全部核心功能',fs=14,c=W,b=True,a=PP_ALIGN.CENTER)

# ====== S17: Summary ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'总结与展望')
Tx(s,Inches(0.8),Inches(1.5),Inches(5.8),Inches(0.4),'项目成果',fs=20,c=D,b=True)
Li(s,Inches(0.8),Inches(2.0),Inches(5.8),Inches(2.5),[
'✓ 完整实现6大核心功能模块',
'✓ 完成2项挑战任务（多态继承+技能系统）',
'✓ 代码量3000+行，15个核心类',
'✓ Win32 GUI中文界面，9按钮8对话框',
'✓ 文件持久化+多线程自动存档',
'✓ 12种道具、4个任务、3种怪物、6种技能',
'✓ 13项测试全部通过，通过率100%',
],fs=14,c=K,sp=Pt(10))
Tx(s,Inches(7.0),Inches(1.5),Inches(5.8),Inches(0.4),'技术收获',fs=20,c=D,b=True)
Li(s,Inches(7.0),Inches(2.0),Inches(5.8),Inches(2.5),[
'继承与多态 — Item/Skill双体系，虚函数/RTTI',
'STL容器 — vector管理对象集合，string/sstream',
'设计模式 — 工厂/原型/策略/模板方法',
'Win32编程 — 消息循环/对话框/GDI字体渲染',
'编码国际化 — UTF-8/GBK/UTF-16转换机制',
'多线程编程 — CreateThread/volatile/同步',
'项目管理 — Git协作/模块化开发/接口约定',
],fs=14,c=K,sp=Pt(10))
Tx(s,Inches(0.8),Inches(5.0),Inches(11.5),Inches(0.4),'未来展望',fs=20,c=D,b=True)
future=[
 ('SQLite数据库','道具/任务/怪物结构化存储','支持复杂查询和管理'),
 ('音效系统','WAV/MP3背景音乐和战斗音效','增强游戏沉浸感'),
 ('地图探索','校园场景地图系统','不同区域分布NPC和怪物'),
 ('装备强化','升级/附魔/套装属性','进阶角色养成玩法'),
 ('网络排行榜','HTTP服务器排名功能','玩家等级和战斗排名'),
 ('跨平台','SDL2/Qt替代Win32 API','支持Linux/macOS'),
]
for i,(title,desc1,desc2) in enumerate(future):
 x=Inches(0.3)+Inches(2.2)*i
 Bx(s,x,Inches(5.6),Inches(2.05),Inches(1.6),G)
 Tx(s,x+Inches(0.05),Inches(5.65),Inches(1.95),Inches(0.3),title,fs=12,c=A,b=True,a=PP_ALIGN.CENTER)
 Tx(s,x+Inches(0.05),Inches(5.98),Inches(1.95),Inches(0.5),desc1,fs=10,c=K,a=PP_ALIGN.CENTER)
 Tx(s,x+Inches(0.05),Inches(6.45),Inches(1.95),Inches(0.5),desc2,fs=9,c=K,a=PP_ALIGN.CENTER)

# ====== S18: Thank You ======
s=P.slides.add_slide(P.slide_layouts[6])
s.background.fill.solid();s.background.fill.fore_color.rgb=D
Bx(s,Inches(0),Inches(0),P.slide_width,Inches(0.12),O)
Bx(s,Inches(0),Inches(7.38),P.slide_width,Inches(0.12),O)
Tx(s,Inches(1),Inches(2.0),Inches(11),Inches(1),'感谢聆听',fs=56,c=W,b=True,a=PP_ALIGN.CENTER)
Bx(s,Inches(5),Inches(3.3),Inches(3.3),Inches(0.03),O)
Tx(s,Inches(1),Inches(3.6),Inches(11),Inches(0.7),'校园RPG冒险游戏 · C++面向对象课程设计',fs=22,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(4.5),Inches(11),Inches(0.5),'组长：张佳怡  |  学号：1030425404',fs=18,c=W,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(5.2),Inches(11),Inches(0.5),'人工智能与自动化学院 · 计算机科学与技术 · 2026年7月',fs=16,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(6.0),Inches(11),Inches(0.5),'欢迎提问与交流',fs=18,c=O,a=PP_ALIGN.CENTER)

OUT = r'C:\Users\张佳怡\OneDrive\Desktop\c++\-\build\校园RPG冒险游戏_汇报.pptx'
P.save(OUT)
print('SUCCESS:'+OUT)
