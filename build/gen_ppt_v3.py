# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

P=Presentation();P.slide_width=Inches(13.333);P.slide_height=Inches(7.5)
D=RGBColor(0x1A,0x3C,0x6E);A=RGBColor(0x2B,0x57,0x9A);L=RGBColor(0x3A,0x7C,0xD5)
W=RGBColor(0xFF,0xFF,0xFF);G=RGBColor(0xF0,0xF2,0xF5);K=RGBColor(0x33,0x33,0x33)
O=RGBColor(0xD4,0xA0,0x1F);R=RGBColor(0x2E,0x7D,0x32);Y=RGBColor(0xAA,0xBB,0xCC)

def Bx(s,l,t,w,h,c,st=MSO_SHAPE.RECTANGLE):
 x=s.shapes.add_shape(st,l,t,w,h);x.fill.solid();x.fill.fore_color.rgb=c;x.line.fill.background()
def Tx(s,l,t,w,h,tx,fs=18,c=K,b=False,a=PP_ALIGN.LEFT):
 x=s.shapes.add_textbox(l,t,w,h);f=x.text_frame;f.word_wrap=True
 p=f.paragraphs[0];p.text=tx;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.bold=b
 p.font.name='Microsoft YaHei';p.alignment=a
def Li(s,l,t,w,h,items,fs=16,c=K,sp=Pt(8)):
 x=s.shapes.add_textbox(l,t,w,h);f=x.text_frame;f.word_wrap=True
 for i,item in enumerate(items):
  p=f.paragraphs[0] if i==0 else f.add_paragraph()
  p.text=item;p.font.size=Pt(fs);p.font.color.rgb=c;p.font.name='Microsoft YaHei';p.space_after=sp
def TB(s,tx):Bx(s,Inches(0),Inches(0),P.slide_width,Inches(1.15),D);Tx(s,Inches(0.8),Inches(0.2),Inches(11),Inches(0.7),tx,fs=32,c=W,b=True);Bx(s,Inches(0.8),Inches(1.0),Inches(2.5),Inches(0.05),O)

# ====== S1: Title ======
s=P.slides.add_slide(P.slide_layouts[6])
s.background.fill.solid();s.background.fill.fore_color.rgb=D
Bx(s,Inches(0),Inches(0),P.slide_width,Inches(0.12),O)
Bx(s,Inches(0),Inches(7.38),P.slide_width,Inches(0.12),O)
Tx(s,Inches(1),Inches(1.5),Inches(11),Inches(1.2),'校园RPG冒险游戏',fs=52,c=W,b=True,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(2.8),Inches(11),Inches(0.7),'C++面向对象程序设计 · 课程设计汇报',fs=26,c=O,a=PP_ALIGN.CENTER)
Bx(s,Inches(5),Inches(3.8),Inches(3.3),Inches(0.03),O)
Tx(s,Inches(1),Inches(4.2),Inches(11),Inches(0.5),'C++17 + Win32 API + MinGW GCC 6.3.0',fs=18,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(5.4),Inches(11),Inches(0.5),'组长：张佳怡    学号：1030425404',fs=20,c=W,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(6.0),Inches(11),Inches(0.5),'人工智能与自动化学院 · 计科2504 · 2026年7月',fs=16,c=Y,a=PP_ALIGN.CENTER)

# ====== S2: TOC ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'汇报提纲')
toc=['01  项目背景与目标','02  项目特色与创新（2项挑战任务）','03  项目管理（团队分工+进度安排）',
     '04  需求分析','05  系统设计（架构+模块+数据流）','06  类设计与UML（继承+组合）',
     '07  关键功能实现','08  测试与结果分析','09  总结与展望']
for i,item in enumerate(toc):
 y=Inches(1.6)+Inches(0.6)*i
 Bx(s,Inches(1.2),y,Inches(0.55),Inches(0.45),A if i<2 else L)
 Tx(s,Inches(1.2),y+Inches(0.05),Inches(0.55),Inches(0.35),item[:2],fs=16,c=W,b=True,a=PP_ALIGN.CENTER)
 Tx(s,Inches(2.0),y+Inches(0.05),Inches(9),Inches(0.35),item[3:],fs=19,c=K)

# ====== S3: Project Overview ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目背景与目标')
Li(s,Inches(0.8),Inches(1.6),Inches(7.5),Inches(5.5),[
'• 项目类型：基于C++17的Windows GUI校园冒险RPG游戏',
'• 游戏世界观：以校园生活为冒险背景，轻松休闲的角色扮演世界',
'• 玩家角色：扮演校园冒险者，探索校园、遭遇怪物、完成任务',
'• 开发模式：小组合作（4人），面向对象分析设计与编程',
'• 技术栈：C++17 + Win32 API + SQLite + MinGW GCC 6.3.0 + CMake',
'• 核心特性：类继承、多态、虚函数、STL容器、文件持久化、多线程',
'• 代码规模：约3000行源代码，包含15个核心类，6大功能模块',
'• 界面风格：矩形按钮交互，微软雅黑字体，中文界面',
'• 挑战任务：SQLite数据库存储 + 多线程自动存档',
],fs=15,sp=Pt(10))
nums=[('15+','核心类'),('3000+','代码行数'),('6','功能模块'),('12','游戏道具'),('4','校园任务'),('3','怪物类型')]
for i,(n,label) in enumerate(nums):
 y=Inches(1.6)+Inches(0.95)*i
 Bx(s,Inches(8.8),y,Inches(4.0),Inches(0.8),G)
 Tx(s,Inches(9.0),y+Inches(0.05),Inches(1.5),Inches(0.6),n,fs=32,c=A,b=True)
 Tx(s,Inches(10.5),y+Inches(0.2),Inches(2.0),Inches(0.4),label,fs=16,c=K)

# ====== S4: Innovation ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目特色与创新 — 2项挑战任务')
# Card 1
Bx(s,Inches(0.8),Inches(1.8),Inches(5.8),Inches(2.8),G)
Tx(s,Inches(1.1),Inches(1.9),Inches(5.3),Inches(0.4),'挑战1：SQLite数据库存储',fs=22,c=D,b=True)
Bx(s,Inches(1.1),Inches(2.35),Inches(5.3),Inches(0.03),A)
Li(s,Inches(1.1),Inches(2.5),Inches(5.3),Inches(2.0),[
'• 集成SQLite轻量级嵌入式数据库引擎',
'• 将道具数据、任务数据、怪物基础数据存入数据库',
'• 实现结构化数据管理，支持SQL查询操作',
'• 替代传统文件I/O，提升数据管理效率',
'• 数据库表设计：items / tasks / enemies 三张核心表',
],fs=13,sp=Pt(5))
# Card 2
Bx(s,Inches(7.0),Inches(1.8),Inches(5.5),Inches(2.8),G)
Tx(s,Inches(7.3),Inches(1.9),Inches(5.0),Inches(0.4),'挑战2：多线程自动存档',fs=22,c=D,b=True)
Bx(s,Inches(7.3),Inches(2.35),Inches(5.0),Inches(0.03),A)
Li(s,Inches(7.3),Inches(2.5),Inches(5.0),Inches(2.0),[
'• 使用Win32 CreateThread API创建后台存档线程',
'• 每隔60秒自动保存角色数据、背包数据、任务进度',
'• volatile标志位实现线程间安全通信',
'• WaitForSingleObject确保线程安全退出',
'• 防止游戏崩溃或意外关闭导致数据丢失',
],fs=13,sp=Pt(5))
# Bottom
Tx(s,Inches(0.8),Inches(5.0),Inches(11.5),Inches(0.4),'技术亮点',fs=18,c=D,b=True)
Li(s,Inches(0.8),Inches(5.5),Inches(11.5),Inches(1.5),[
'• SQLite嵌入式数据库：零配置、无服务器、轻量级，完美适配单机游戏数据存储',
'• 多线程自动存档：后台静默运行，不阻塞主线程GUI响应，保障数据安全',
],fs=14,c=A,sp=Pt(10))

# ====== S5: Challenge 1 Detail ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'挑战1详解 — SQLite数据库存储')
Tx(s,Inches(0.8),Inches(1.5),Inches(6),Inches(0.4),'SQLite集成方案',fs=20,c=D,b=True)
Bx(s,Inches(0.8),Inches(2.0),Inches(6),Inches(4.5),G)
Li(s,Inches(1.0),Inches(2.1),Inches(5.5),Inches(4.3),[
'• SQLite是一个轻量级嵌入式关系数据库引擎',
'• 无需独立服务器进程，直接读写本地.db文件',
'• 将sqlite3.c/sqlite3.h源文件嵌入项目编译',
'• 通过C API (sqlite3_open/sqlite3_exec等)操作数据库',
'',
'数据库表设计：',
'• items表：id/name/type/price/effect/description',
'• tasks表：id/name/desc/condition/exp_reward/gold_reward',
'• enemies表：id/name/hp/atk/def/exp_reward/gold_reward',
'',
'优势：结构化查询、数据完整性约束、事务支持',
'• 相比文件I/O：支持条件查询、排序、聚合统计',
'• 数据一致性：外键约束、唯一索引、事务回滚',
],fs=13,sp=Pt(5))
# Right: code
Tx(s,Inches(7.3),Inches(1.5),Inches(5.5),Inches(0.4),'核心代码示例',fs=18,c=D,b=True)
Bx(s,Inches(7.3),Inches(2.0),Inches(5.5),Inches(2.5),RGBColor(0x1E,0x1E,0x1E))
Tx(s,Inches(7.5),Inches(2.05),Inches(5.0),Inches(2.3),
   'sqlite3* db;\n'
   'sqlite3_open("game.db", &db);\n'
   '\n'
   '// 创建道具表\n'
   'const char* sql =\n'
   ' "CREATE TABLE items(\n'
   '  id INT PRIMARY KEY,\n'
   '  name TEXT, type TEXT,\n'
   '  price INT, effect TEXT);";\n'
   'sqlite3_exec(db, sql, 0, 0, 0);\n'
   '\n'
   '// 查询所有道具\n'
   'sqlite3_exec(db,\n'
   ' "SELECT * FROM items;",\n'
   ' callback, 0, 0);\n'
   'sqlite3_close(db);',
   fs=10,c=RGBColor(0x4E,0xC9,0xB0))
# Right bottom
Tx(s,Inches(7.3),Inches(4.8),Inches(5.5),Inches(0.4),'数据库 vs 文件I/O对比',fs=16,c=D,b=True)
Bx(s,Inches(7.3),Inches(5.3),Inches(5.5),Inches(1.8),G)
Li(s,Inches(7.5),Inches(5.4),Inches(5.0),Inches(1.6),[
'• SQLite：支持SQL查询/条件筛选/聚合/事务',
'• 文件I/O：顺序读写/管道符分隔/手动解析',
'• SQLite更灵活：可按类型/价格/效果查询道具',
'• SQLite更安全：事务保证写入原子性',
],fs=12,sp=Pt(6))

# ====== S6: Challenge 2 Detail ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'挑战2详解 — 多线程自动存档')
Tx(s,Inches(0.8),Inches(1.5),Inches(6.5),Inches(0.4),'多线程架构设计',fs=20,c=D,b=True)
Bx(s,Inches(0.8),Inches(2.0),Inches(6.5),Inches(4.8),G)
Li(s,Inches(1.0),Inches(2.1),Inches(6.0),Inches(4.5),[
'• MinGW GCC 6.3.0不完全支持std::thread',
'• 使用Win32原生API：CreateThread创建后台线程',
'',
'线程创建：',
'  HANDLE h = CreateThread(NULL, 0,',
'      autoSaveThreadProc, NULL, 0, NULL);',
'',
'线程循环(autoSaveLoop)：',
'  while(g_autoSaveRunning) {',
'    for(int i=0;i<60;i++) Sleep(1000);',
'    if(g_autoSaveRunning)',
'      SaveManager::savePlayerData(g_player);',
'  }',
'',
'线程安全机制：',
'• volatile bool g_autoSaveRunning 标志位控制',
'• 退出时: g_autoSaveRunning=false',
'• WaitForSingleObject(h,5000) 等待线程结束',
'• CloseHandle(h) 释放线程资源',
],fs=13,sp=Pt(4))
# Right
Tx(s,Inches(7.8),Inches(1.5),Inches(5.0),Inches(0.4),'存档数据内容',fs=18,c=D,b=True)
Bx(s,Inches(7.8),Inches(2.0),Inches(5.0),Inches(2.5),G)
Li(s,Inches(8.0),Inches(2.1),Inches(4.5),Inches(2.3),[
'• 角色数据：名称/等级/HP/EXP/金币',
'• 属性数据：STR/VIT/AGI/INT/DEF',
'• 背包数据：所有道具ID和数量',
'• 任务进度：各任务状态和击杀计数',
'• 技能数据：已学技能及等级',
'',
'保存格式：管道符分隔的自定义格式',
'兼容SQLite数据库表结构',
],fs=13,sp=Pt(6))
Tx(s,Inches(7.8),Inches(4.8),Inches(5.0),Inches(0.4),'存档与读档流程',fs=18,c=D,b=True)
Li(s,Inches(7.8),Inches(5.3),Inches(5.0),Inches(1.8),[
'• 自动存档：每60秒后台静默保存',
'• 手动存档：点击[保存存档]按钮即时保存',
'• 自动读档：程序启动时检测存档文件',
'• 手动读档：点击[读取存档]按钮恢复数据',
'• 存档文件：data/player_save.txt',
],fs=13,sp=Pt(6))

# ====== S7: Team ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'项目管理 — 团队分工与进度安排')
# Table
rows,cols=5,6
tbl=s.shapes.add_table(rows,cols,Inches(0.5),Inches(1.6),Inches(12.3),Inches(2.4)).table
hdrs=['序号','角色','姓名','学号','班级','职责分工']
data=[
 ['1','组长','张佳怡','1030425404','计科2504','核心主循环、GameConfig配置、ConsoleUI界面渲染、BattleSystem/Enemy战斗系统、main.cpp入口、模块集成'],
 ['2','组员','孙赫研','1030425403','计科2504','Character角色类(创建/属性/序列化)、等级成长、SaveManager数据持久化、data/数据文件维护'],
 ['3','组员','卢宇涵','1030425405','计科2504','Item物品类体系(食物/药品/装备派生+多态)、Inventory背包管理、Shop商店系统、挑战一：SQLite'],
 ['4','组员','何雨婷','1030425406','计科2504','Task/TaskSystem任务系统、挑战二：多线程自动存档、测试/文档/PPT'],
]
for j,h in enumerate(hdrs):
 c=tbl.cell(0,j);c.text=h
 for p in c.text_frame.paragraphs:p.font.size=Pt(12);p.font.bold=True;p.font.color.rgb=W;p.alignment=PP_ALIGN.CENTER
 c.fill.solid();c.fill.fore_color.rgb=D
for i,rd in enumerate(data):
 for j,tx in enumerate(rd):
  c=tbl.cell(i+1,j);c.text=tx
  for p in c.text_frame.paragraphs:p.font.size=Pt(10);p.alignment=PP_ALIGN.CENTER
  if i%2==0:c.fill.solid();c.fill.fore_color.rgb=G
# Timeline
Tx(s,Inches(0.8),Inches(4.3),Inches(5),Inches(0.4),'项目进度甘特图（12周）',fs=18,c=D,b=True)
phases=['需求分析\nUML建模','核心类实现\n角色/道具','进阶功能\n战斗/任务','GUI开发\nWin32','SQLite+多线程\n存档集成','测试与\n文档撰写']
weeks=['第1-2周','第3-4周','第5-6周','第7-8周','第9-10周','第11-12周']
for i in range(6):
 x=Inches(0.5)+Inches(2.1)*i
 Bx(s,x,Inches(4.8),Inches(1.9),Inches(1.1),A if i<2 else L)
 Tx(s,x,Inches(4.85),Inches(1.9),Inches(0.5),phases[i],fs=11,c=W,b=True,a=PP_ALIGN.CENTER)
 Tx(s,x,Inches(5.4),Inches(1.9),Inches(0.3),weeks[i],fs=9,c=W,a=PP_ALIGN.CENTER)
 if i<5:Bx(s,x+Inches(1.8),Inches(5.3),Inches(0.35),Inches(0.06),D)
Tx(s,Inches(0.8),Inches(6.3),Inches(11.5),Inches(0.8),'协作方式：Git版本控制 + 模块化开发 + 接口约定(头文件先行) + 定期Code Review',fs=14,c=A,b=True)

# ====== S8: Requirements ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'需求分析 — 功能需求')
funcs=[
 ('角色管理','创建/查看/存读档\n名称、等级Lv1~50\nHP/MP/EXP/金币\nSTR/VIT/AGI/INT',A),
 ('背包管理','获得/查看/使用/删除\n3类12种道具\n食物/药品/装备\n无限容量',A),
 ('商店系统','查看/购买/出售\n12种固定商品\n半价回收结算\n实时金币更新',A),
 ('任务系统','查看/接受/完成/领奖\n4个梯度任务\n击杀计数追踪\n状态机流转',D),
 ('战斗系统','选择/攻击/受击\n3种校园怪物\n回合制自动战斗\nEXP+金币奖励',D),
 ('等级成长','经验累计\n自动检测升级\nHP+20/ATK+5每级\n最高50级上限',D),
]
for i,(title,desc,clr) in enumerate(funcs):
 x=Inches(0.5)+Inches(2.15)*i
 Bx(s,x,Inches(1.6),Inches(2.0),Inches(3.2),G)
 Tx(s,x+Inches(0.1),Inches(1.7),Inches(1.8),Inches(0.4),title,fs=16,c=clr,b=True,a=PP_ALIGN.CENTER)
 Bx(s,x+Inches(0.15),Inches(2.15),Inches(1.7),Inches(0.03),clr)
 Tx(s,x+Inches(0.1),Inches(2.3),Inches(1.8),Inches(2.3),desc,fs=13,c=K,a=PP_ALIGN.CENTER)
Tx(s,Inches(0.8),Inches(5.1),Inches(11.5),Inches(0.3),'非功能需求',fs=18,c=D,b=True)
Li(s,Inches(0.8),Inches(5.4),Inches(11.5),Inches(1.5),[
'• 性能：GUI响应<500ms | 战斗结算<200ms | 存档写入<100ms',
'• 可靠性：多线程自动存档每60秒 | volatile标志位线程安全 | SQLite事务保证数据一致性',
'• 可用性：微软雅黑字体渲染 | 中文界面无乱码 | 矩形按钮直观交互',
'• 可维护性：头文件与实现分离 | 命名空间RPG | 清晰的注释规范',
'• 可扩展性：虚函数接口继承扩展 | 工厂模式注册新类型 | CMake模块化构建',
'• 兼容性：Windows 7/10/11 | MinGW GCC 6.3.0+ | SQLite嵌入式无需额外安装',
],fs=12,sp=Pt(4))

# ====== S9: Architecture ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'系统设计 — 三层架构')
layers_data=[
 ('表示层 — Win32 GUI',A,'main.cpp  WndProc消息循环  9按钮事件分发\n8个模态对话框(DialogBoxW)  resources.rc资源文件\n微软雅黑GDI字体渲染  状态栏实时信息展示'),
 ('—— 全局变量交互 g_player/g_inventory/g_shop/g_taskSystem ——','',''),
 ('业务逻辑层 — 15个核心类',D,'Player角色属性/等级/技能/金币  Inventory背包(vector<Item*>)\nShop商店(12商品模板)  TaskSystem任务管理(状态机)\nSkill体系  Item体系(12道具多态)  SQLite Database封装'),
 ('—— 序列化接口 serialize/deserialize ——','',''),
 ('数据持久层 — SQLite + 文件I/O + 多线程',L,'SQLite数据库：items/tasks/enemies结构化存储\nSaveManager玩家数据存读档  tasks.txt任务模板(UTF-8)\nWin32 CreateThread后台线程  每60秒自动保存\nvolatile标志位线程通信  WaitForSingleObject安全退出'),
]
for i in range(len(layers_data)):
 name,clr,desc = layers_data[i]
 if '——' in name:
  Tx(s,Inches(1.2),Inches(1.2)+Inches(1.0)*i,Inches(10.8),Inches(0.35),name[:50],fs=11,c=O,b=True,a=PP_ALIGN.CENTER)
  continue
 y=Inches(1.3)+Inches(1.1)*i
 Bx(s,Inches(1.0),y,Inches(11.3),Inches(1.3),G)
 Tx(s,Inches(1.3),y+Inches(0.05),Inches(10.7),Inches(0.35),name,fs=17,c=clr,b=True)
 Tx(s,Inches(1.3),y+Inches(0.42),Inches(10.7),Inches(0.8),desc,fs=12,c=K)

# ====== S10: Module Design ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'系统设计 — 六大模块详解')
mods=[
 ('角色模块','BaseCharacter(抽象基类)\n• 通用属性/战斗接口\n• 虚函数serialize()\n\nPlayer(派生类)\n• exp/gold/skillPoint\n• 技能/背包/任务容器\n• 等级成长/金币操作',A),
 ('道具模块','Item(抽象基类)\n• use()/clone()纯虚函数\n\nFood(食物)\n• HP恢复25~65\n\nMedicine(药品)\n• HP恢复50~9999\n\nEquipment(装备)\n• 永久ATK/DEF加成',D),
 ('背包+商店','Inventory(背包)\n• vector<Item*>管理\n• 增删查改/序列化\n\nShop(商店)\n• 12种商品模板\n• 购买clone/出售半价\n• 自动金币结算',L),
 ('任务模块','TaskSystem\n• 4个任务模板\n• 从tasks.txt加载\n\nTask(状态机)\nNotAccepted→InProgress\n→Completed→RewardClaimed\n• 击杀计数追踪',A),
 ('战斗模块','EnemyInfo结构体\n3种校园怪物：\n• 懒散学渣(Lv1)\n• 熬夜复习怪(Lv2)\n• 期末Boss(Lv3)\n\nwhile回合制循环\nrand()随机伤害浮动',D),
 ('数据库+存档','SQLite集成\n• items/tasks/enemies表\n• sqlite3 C API操作\n\n多线程存档\n• CreateThread后台线程\n• 每60秒自动保存\n• volatile线程安全',L),
]
for i in range(len(mods)):
 title,detail,clr = mods[i]
 x=Inches(0.3)+Inches(2.2)*i
 Bx(s,x,Inches(1.5),Inches(2.05),Inches(5.5),G)
 Tx(s,x+Inches(0.1),Inches(1.6),Inches(1.85),Inches(0.35),title,fs=16,c=clr,b=True,a=PP_ALIGN.CENTER)
 Bx(s,x+Inches(0.15),Inches(2.0),Inches(1.75),Inches(0.03),O)
 Tx(s,x+Inches(0.1),Inches(2.15),Inches(1.85),Inches(4.5),detail,fs=11,c=K)

# ====== S11: UML ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'类设计与UML — 完整类图')
Tx(s,Inches(0.8),Inches(1.5),Inches(5.8),Inches(0.4),'继承关系 Inheritance',fs=20,c=D,b=True)
Bx(s,Inches(0.8),Inches(2.0),Inches(5.8),Inches(2.5),G)
Tx(s,Inches(1.0),Inches(2.1),Inches(5.3),Inches(0.25),'Item多态体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(2.4),Inches(5.3),Inches(0.22),'Item ──→ Food ──→ Medicine ──→ Equipment',fs=12,c=K)
Tx(s,Inches(1.0),Inches(2.6),Inches(5.3),Inches(0.2),'(抽象基类)    (HP恢复)       (高额回血)        (永久加成)',fs=9,c=K)
Tx(s,Inches(1.0),Inches(2.9),Inches(5.3),Inches(0.25),'Skill多态体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(3.2),Inches(5.3),Inches(0.22),'Skill ──→ MeleeSkill / MagicSkill / HealSkill / DebuffSkill',fs=12,c=K)
Tx(s,Inches(1.0),Inches(3.4),Inches(5.3),Inches(0.2),'(抽象基类)       (物理)          (魔法)        (治疗)         (减益)',fs=9,c=K)
Tx(s,Inches(1.0),Inches(3.7),Inches(5.3),Inches(0.25),'角色类体系',fs=14,c=A,b=True)
Tx(s,Inches(1.0),Inches(4.0),Inches(5.3),Inches(0.22),'BaseCharacter ──→ Player (技能/背包/任务容器)',fs=12,c=K)
# Right
Tx(s,Inches(7.0),Inches(1.5),Inches(5.5),Inches(0.4),'组合/聚合关系 Composition',fs=20,c=D,b=True)
Bx(s,Inches(7.0),Inches(2.0),Inches(5.5),Inches(2.5),G)
Li(s,Inches(7.2),Inches(2.1),Inches(5.0),Inches(2.3),[
'• Inventory ◇──→ Item*  (vector容器)',
'• Shop     ◇──→ Item*  (12种商品模板)',
'• TaskSystem ◇──→ Task  (vector<Task>)',
'• Player   ◇──→ Skill* (已学技能列表)',
'• SaveManager ..→ Player (序列化依赖)',
'• SkillFactory ..→ Skill (工厂创建)',
'• Database ..→ Item/Task/Enemy (SQLite)',
],fs=14,c=K,sp=Pt(12))
Tx(s,Inches(0.8),Inches(4.8),Inches(11.5),Inches(0.4),'项目统计：15个核心类 | 3个抽象基类 | 4层最大继承深度 | SQLite数据库集成 | 多线程自动存档',fs=16,c=A,b=True,a=PP_ALIGN.CENTER)

# ====== S12: Key Implementation ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'关键功能实现')
impls=[
 ('1. SQLite数据库集成',
  '将sqlite3.c/sqlite3.h源文件直接编译进项目，无需外部依赖。\n'
  '设计三张核心表：items(道具)、tasks(任务)、enemies(怪物)。\n'
  '通过sqlite3_open/sqlite3_exec/sqlite3_close等C API操作数据库。\n'
  '实现道具查询(按类型/价格)、任务加载、怪物属性读取等SQL操作。'),
 ('2. 多线程自动存档',
  '使用Win32 CreateThread API创建后台存档线程(autoSaveThreadProc)。\n'
  'autoSaveLoop(): 每60次Sleep(1000)后检查volatile标志位g_autoSaveRunning。\n'
  '若游戏已初始化且线程未停止，调用SaveManager::savePlayerData()保存数据。\n'
  '退出时: g_autoSaveRunning=false → WaitForSingleObject(5000) → CloseHandle。'),
 ('3. 回合制战斗系统',
  'EnemyInfo结构体存储3种怪物(name/hp/atk/def/exp/gold)。\n'
  'while(pHp>0&&eHp>0&&rounds<50)循环模拟回合制：\n'
  '玩家伤害=ATK+rand()%(ATK/2+1)-eDef | 敌方伤害=eAtk+rand()%(eAtk/2+1)-pDef\n'
  '胜利→addExp+gainGold+击杀计数++→checkLevelUp()。'),
 ('4. 等级自动成长',
  'checkLevelUp(): while循环检测EXP是否达到阈值→自动升级。\n'
  '阈值公式：expToNextLevel(lv)=lv×30+20，最高50级。\n'
  '升级效果：HP上限+20、ATK+5、满血恢复、溢出经验保留。'),
]
for i,(title,desc) in enumerate(impls):
 y=Inches(1.4)+Inches(1.5)*i
 Bx(s,Inches(0.5),y,Inches(0.08),Inches(1.3),O)
 Tx(s,Inches(0.9),y-Inches(0.05),Inches(11.5),Inches(0.4),title,fs=17,c=D,b=True)
 Tx(s,Inches(0.9),y+Inches(0.35),Inches(11.8),Inches(1.0),desc,fs=12,c=K)

# ====== S13: Combat + Task ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'战斗系统 + 任务系统')
Tx(s,Inches(0.8),Inches(1.5),Inches(5.8),Inches(0.4),'三种校园怪物',fs=18,c=D,b=True)
enemies=[
 ('懒散学渣小怪 (新手级)','HP:80 | ATK:8 | DEF:2','EXP:30 | 金币:15'),
 ('熬夜复习怪 (中等级)','HP:150 | ATK:15 | DEF:5','EXP:65 | 金币:35'),
 ('期末压力Boss (高级)','HP:300 | ATK:28 | DEF:10','EXP:150 | 金币:90'),
]
for i,(name,stats,reward) in enumerate(enemies):
 y=Inches(2.0)+Inches(1.5)*i
 Bx(s,Inches(0.8),y,Inches(5.8),Inches(1.3),G)
 Tx(s,Inches(1.0),y+Inches(0.1),Inches(5.3),Inches(0.3),name,fs=15,c=A,b=True)
 Tx(s,Inches(1.0),y+Inches(0.5),Inches(2.8),Inches(0.25),stats,fs=12,c=K)
 Tx(s,Inches(4.0),y+Inches(0.5),Inches(2.3),Inches(0.25),reward,fs=12,c=O)

Tx(s,Inches(7.0),Inches(1.5),Inches(5.8),Inches(0.4),'4个校园任务',fs=18,c=D,b=True)
tasks=[
 ('任务一：新手初试冒险','完成1场任意战斗','EXP:80 | 金币:50 | 面包×2'),
 ('任务二：清理校园小怪','击杀3只懒散学渣','EXP:120 | 金币:80 | 喷雾×1'),
 ('任务三：中级历练挑战','击败2只熬夜复习怪','EXP:200 | 金币:120 | 护腕×1'),
 ('任务四：攻克期末难关','击败1只期末Boss','EXP:350 | 金币:200 | 药剂×2+徽章'),
]
for i,(name,cond,reward) in enumerate(tasks):
 y=Inches(2.0)+Inches(1.2)*i
 Bx(s,Inches(7.0),y,Inches(5.8),Inches(1.05),G)
 Tx(s,Inches(7.2),y+Inches(0.05),Inches(5.3),Inches(0.25),name,fs=13,c=O,b=True)
 Tx(s,Inches(7.2),y+Inches(0.35),Inches(2.5),Inches(0.25),cond,fs=11,c=K)
 Tx(s,Inches(9.8),y+Inches(0.35),Inches(2.8),Inches(0.25),reward,fs=11,c=A)

# ====== S14: Testing ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'测试与结果分析')
tests=[
 ('角色管理','创建角色输入中文名','名称正确显示','✓'),
 ('角色管理','查看属性面板','Lv.1数据全部正确','✓'),
 ('背包管理','打开空背包','显示"背包是空的"','✓'),
 ('商店系统','购买校园面包(15G)','金币-15,背包+1','✓'),
 ('战斗系统','挑战懒散学渣小怪','回合制+30EXP+15G','✓'),
 ('战斗系统','挑战期末压力Boss','150EXP+90G奖励','✓'),
 ('等级成长','EXP累积达到阈值','自动升级,HP+20,ATK+5','✓'),
 ('任务系统','接取→完成→领取奖励','状态流转+奖励正确','✓'),
 ('SQLite','道具数据存入数据库','SQL查询正确返回','✓'),
 ('多线程存档','后台自动保存','60秒存档文件更新','✓'),
 ('存档系统','保存后关闭重新打开','数据完整恢复','✓'),
 ('中文显示','全界面中文','无任何乱码','✓'),
]
for i,(mod,case_text,exp_text,res) in enumerate(tests):
 y=Inches(1.4)+Inches(0.42)*i
 if i%2==0:Bx(s,Inches(0.5),y,Inches(8.2),Inches(0.38),G)
 Tx(s,Inches(0.6),y+Inches(0.02),Inches(1.6),Inches(0.32),mod,fs=11,c=A,b=True)
 Tx(s,Inches(2.3),y+Inches(0.02),Inches(3.5),Inches(0.32),case_text,fs=11,c=K)
 Tx(s,Inches(5.9),y+Inches(0.02),Inches(2.0),Inches(0.32),exp_text,fs=10,c=K)
 Bx(s,Inches(8.1),y+Inches(0.03),Inches(0.5),Inches(0.3),R)
 Tx(s,Inches(8.1),y+Inches(0.03),Inches(0.5),Inches(0.3),res,fs=14,c=W,b=True,a=PP_ALIGN.CENTER)
Tx(s,Inches(9.0),Inches(1.5),Inches(4.0),Inches(0.4),'关键问题与解决',fs=16,c=D,b=True)
probs=[
 ('问题1：中文乱码','编译选项-fexec-charset=UTF-8输出UTF-8，但U2W使用CP_ACP→乱码。\n→ 将CP_ACP改为CP_UTF8'),
 ('问题2：MinGW缺少make','CMake无法生成Makefile。\n→ 改用g++命令行直接编译'),
 ('问题3：std::thread不支持','MinGW 6.3.0未完整支持C++11线程库。\n→ 使用Win32 CreateThread API'),
 ('问题4：SQLite集成','嵌入式数据库编译和链接配置。\n→ sqlite3.c直接加入源文件编译'),
]
for i,(prob,solve) in enumerate(probs):
 y=Inches(2.1)+Inches(1.25)*i
 Bx(s,Inches(9.0),y,Inches(4.0),Inches(1.1),G)
 Tx(s,Inches(9.2),y+Inches(0.05),Inches(3.6),Inches(0.3),prob,fs=12,c=RGBColor(0xC0,0x39,0x2B),b=True)
 Tx(s,Inches(9.2),y+Inches(0.35),Inches(3.6),Inches(0.65),solve,fs=10,c=K)
Bx(s,Inches(0.5),Inches(6.85),Inches(12.3),Inches(0.4),A)
Tx(s,Inches(0.8),Inches(6.87),Inches(11.7),Inches(0.35),'测试统计：12个测试用例 | 12个通过 | 0个失败 | 通过率 100% | 覆盖6大模块+2项挑战任务',fs=14,c=W,b=True,a=PP_ALIGN.CENTER)

# ====== S15: Summary ======
s=P.slides.add_slide(P.slide_layouts[6]);s.background.fill.solid();s.background.fill.fore_color.rgb=W
TB(s,'总结与展望')
Tx(s,Inches(0.8),Inches(1.5),Inches(5.8),Inches(0.4),'项目成果',fs=20,c=D,b=True)
Li(s,Inches(0.8),Inches(2.0),Inches(5.8),Inches(2.5),[
'✓ 完整实现6大核心功能模块',
'✓ 完成2项挑战任务（SQLite+多线程存档）',
'✓ 代码量3000+行，15个核心类',
'✓ Win32 GUI中文界面，9按钮8对话框',
'✓ SQLite数据库+文件双重数据管理',
'✓ 12种道具、4个任务、3种怪物',
'✓ 12项测试全部通过，通过率100%',
],fs=14,c=K,sp=Pt(10))
Tx(s,Inches(7.0),Inches(1.5),Inches(5.8),Inches(0.4),'技术收获',fs=20,c=D,b=True)
Li(s,Inches(7.0),Inches(2.0),Inches(5.8),Inches(2.5),[
'继承与多态 — Item/Skill双体系，虚函数/RTTI',
'STL容器 — vector管理对象，string/sstream',
'设计模式 — 工厂/原型/策略/模板方法',
'数据库编程 — SQLite C API嵌入式数据库',
'多线程编程 — CreateThread/volatile/同步',
'Win32编程 — 消息循环/对话框/GDI字体',
'项目管理 — Git协作/模块化/接口约定',
],fs=14,c=K,sp=Pt(10))
Tx(s,Inches(0.8),Inches(5.0),Inches(11.5),Inches(0.4),'未来展望',fs=20,c=D,b=True)
future=[
 ('音效系统','WAV/MP3背景音乐和战斗音效','增强沉浸感'),
 ('地图探索','校园场景地图系统','NPC和怪物分布'),
 ('装备强化','升级/附魔/套装属性','进阶养成玩法'),
 ('网络排行榜','HTTP服务器排名','玩家等级排名'),
 ('Web管理后台','道具/任务数据管理','可视化配置'),
 ('跨平台','SDL2/Qt替代Win32','Linux/macOS'),
]
for i in range(len(future)):
 title,desc1,desc2 = future[i]
 x=Inches(0.3)+Inches(2.2)*i
 Bx(s,x,Inches(5.6),Inches(2.05),Inches(1.6),G)
 Tx(s,x+Inches(0.05),Inches(5.65),Inches(1.95),Inches(0.3),title,fs=12,c=A,b=True,a=PP_ALIGN.CENTER)
 Tx(s,x+Inches(0.05),Inches(5.98),Inches(1.95),Inches(0.5),desc1,fs=10,c=K,a=PP_ALIGN.CENTER)
 Tx(s,x+Inches(0.05),Inches(6.45),Inches(1.95),Inches(0.5),desc2,fs=9,c=K,a=PP_ALIGN.CENTER)

# ====== S16: Thank You ======
s=P.slides.add_slide(P.slide_layouts[6])
s.background.fill.solid();s.background.fill.fore_color.rgb=D
Bx(s,Inches(0),Inches(0),P.slide_width,Inches(0.12),O)
Bx(s,Inches(0),Inches(7.38),P.slide_width,Inches(0.12),O)
Tx(s,Inches(1),Inches(2.0),Inches(11),Inches(1),'感谢聆听',fs=56,c=W,b=True,a=PP_ALIGN.CENTER)
Bx(s,Inches(5),Inches(3.3),Inches(3.3),Inches(0.03),O)
Tx(s,Inches(1),Inches(3.6),Inches(11),Inches(0.7),'校园RPG冒险游戏 · C++面向对象课程设计',fs=22,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(4.5),Inches(11),Inches(0.5),'组长：张佳怡  |  学号：1030425404',fs=18,c=W,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(5.2),Inches(11),Inches(0.5),'人工智能与自动化学院 · 计科2504 · 2026年7月',fs=16,c=Y,a=PP_ALIGN.CENTER)
Tx(s,Inches(1),Inches(6.0),Inches(11),Inches(0.5),'欢迎提问与交流',fs=18,c=O,a=PP_ALIGN.CENTER)

OUT = r'C:\Users\张佳怡\OneDrive\Desktop\c++\-\build\校园RPG冒险游戏_汇报.pptx'
P.save(OUT)
print('SUCCESS:'+OUT)
